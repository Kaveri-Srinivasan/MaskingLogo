import os
import shutil
import hashlib
import difflib
import docx
import pptx
import PyPDF2
import pandas as pd
from collections import defaultdict

def get_text_from_file(file_path):
    """Extracts text from different file types."""
    ext = file_path.lower().split('.')[-1]
    text = ""
    
    if ext == 'txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
    elif ext == 'pdf':
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text
        except Exception as e:
            print(f"Error processing {file_path}: {e}")

    # elif ext == 'pdf':
    #     with open(file_path, 'rb') as f:
    #         reader = PyPDF2.PdfReader(f)
    #         for page in reader.pages:
    #             text += page.extract_text() or ""
    elif ext == 'docx':
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif ext in ['xls', 'xlsx']:
        df = pd.read_excel(file_path, engine='openpyxl')
        text = df.to_string()
    elif ext in ['ppt', 'pptx']:
        presentation = pptx.Presentation(file_path)
        text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
    
    
    return text.strip()

def similarity_score(text1, text2):
    """Returns similarity ratio between two texts."""
    return difflib.SequenceMatcher(None, text1, text2).ratio()

def process_documents(input_folder, output_folder, threshold=0.85):
    """Processes documents, groups duplicates, and moves unique files."""
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    duplicate_folder = os.path.join(output_folder, 'duplicates')
    unique_folder = os.path.join(output_folder, 'unique')
    os.makedirs(duplicate_folder, exist_ok=True)
    os.makedirs(unique_folder, exist_ok=True)
    
    files = [os.path.join(input_folder, f) for f in os.listdir(input_folder) if os.path.isfile(os.path.join(input_folder, f))]
    file_texts = {}
    grouped_duplicates = defaultdict(list)
    
    for file in files:
        text = get_text_from_file(file)
        if text:
            file_texts[file] = text
    
    processed = set()
    for file1, text1 in file_texts.items():
        if file1 in processed:
            continue
        duplicates = [file1]
        for file2, text2 in file_texts.items():
            if file1 != file2 and similarity_score(text1, text2) >= threshold:
                duplicates.append(file2)
                processed.add(file2)
        
        if len(duplicates) > 1:
            folder_name = hashlib.md5(text1.encode()).hexdigest()
            duplicate_subfolder = os.path.join(duplicate_folder, folder_name)
            os.makedirs(duplicate_subfolder, exist_ok=True)
            for dup_file in duplicates:
                shutil.move(dup_file, os.path.join(duplicate_subfolder, os.path.basename(dup_file)))
        else:
            shutil.move(file1, os.path.join(unique_folder, os.path.basename(file1)))

def main():
    input_folder = r"C:\Users\kaveri.s\Downloads\OneDrive_2_3-25-2025 (1) 1"
    output_folder = r"C:\Users\kaveri.s\Downloads"
  # Change this to your destination folder
    threshold = 0.85  # Duplicate threshold (85%)
    
    process_documents(input_folder, output_folder, threshold)
    print("Processing complete! Check the output folder.")

if __name__ == "__main__":
    main()
